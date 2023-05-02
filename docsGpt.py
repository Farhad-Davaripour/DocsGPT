# docsGpt.py - Contains the docsGpt functions and classes for document parsing
# Author: Armin Norouzi, Farhad Davaripour 
# Contact: https://github.com/Farhad-Davaripour/DocsGPT
# Date created: April 14, 2023
# Last modified: May 3, 2023
# License: MIT License

# Import required modules
import sys
import subprocess
from google.colab import files
import os
import shutil
import time
import tempfile

# List of library names to import
library_names = ['langchain', 'openai', 'PyPDF2', 'tiktoken', 'faiss-cpu', 'textwrap', 'python-docx', 'python-pptx']

# Dynamically import libraries from list
for name in library_names:
    try:
        __import__(name)
    except ImportError:
        print(f"{name} not found. Installing {name}...")
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', name])


# Import required modules
from PyPDF2 import PdfReader 
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import ElasticVectorSearch, Pinecone, Weaviate, FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from getpass import getpass
import textwrap
import os
import docx
import pptx

# adding token
# print("You need OpenAI token: Here is the link to get 
# the keys: https://platform.openai.com/account/billing/overview")

token = getpass("Enter your OpenAI token: ()")
os.environ["OPENAI_API_KEY"] = str(token)


# Download embeddings from OpenAI
embeddings = OpenAIEmbeddings()
chain = load_qa_chain(OpenAI(), chain_type="stuff")


def extract_texts(root_files):
    """
    Extracts text from uploaded file and puts it in a list.
    Supported file types: .pdf, .docx, .pptx
    If multiple files are provided, their contents are concatenated.
    Args:
    - root_files: A list of file paths to be processed.
    Returns:
    - A FAISS index object containing the embeddings of the 
    text chunks.
    """
    raw_text = ''

    for root_file in root_files:
        _, ext = os.path.splitext(root_file)
        if ext == '.pdf':
            with open(root_file, 'rb') as f:
                reader = PdfReader(f)
                for i in range(len(reader.pages)):
                    page = reader.pages[i]
                    raw_text += page.extract_text()
        elif ext == '.docx':
            doc = docx.Document(root_file)
            for paragraph in doc.paragraphs:
                raw_text += paragraph.text
        elif ext == '.pptx':
            ppt = pptx.Presentation(root_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        raw_text += shape.text

    # retreival we don't hit the token size limits. 
    text_splitter = CharacterTextSplitter(        
                                            separator = "\n",
                                            chunk_size = 1000,
                                            chunk_overlap  = 200,
                                            length_function = len,
                                        )

    texts = text_splitter.split_text(raw_text)

    docsearch = FAISS.from_texts(texts, embeddings)
    return docsearch


def run_query(query, docsearch):
    """
    Runs a query on a PDF file using the docsearch and chain 
    libraries. 
    Args:
    - query: A string representing the query to be run.
    - file: A PDFReader object containing the PDF file to be 
    searched.
    Returns:
    - A string containing the output of the chain library run 
    on the documents returned by the docsearch similarity search.
    """
    
    docs = docsearch.similarity_search(query)
    return chain.run(input_documents=docs, question=query)


def upload_file(folder_path):
    """
    Uploads a file from the local file system and saves it to 
    a folder path. 
    Args:
    - folder_path: A string representing the folder path where 
    the file will be saved.
    Returns:
    - A string representing the path of the uploaded file.
    """
    
    uploaded = files.upload()
    root_file = []

    for filename, data in uploaded.items():
        with open(filename, 'wb') as f:
            f.write(data)
        shutil.copy(filename, folder_path + "/")
        root_file.append(folder_path + "/" + filename)
        os.remove(filename)


    return root_file


def run_conversation(folder_path):
    """
    Initiates a conversation with the user by repeatedly asking for 
    input queries and running them on a PDF file. 
    Args:
    - folder_path: A string representing the folder path where the 
    PDF file is located.
    Returns:
    - Run conversation based on PDF
    """
    root_files = upload_file(folder_path)
    # location of the pdf file/files.


    docsearch = extract_texts(root_files)

    count = 0
    while True:
        print("Question ", count + 1)

        query = input(" Ask your question or if you have no further question type stop:\n ")
        
        if query.lower() == "stop":
            print("### Thanks for using the app! ###")
            break
        elif query == "":
            print("### Your input is empty! Try again! ###")
            continue
        else:
            wrapped_text = textwrap.wrap(run_query(query, docsearch), width=100)
            print("Answer:")
            for line in wrapped_text:
                print(line)
            count += 1